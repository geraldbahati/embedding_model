from docx import Document
from pptx import Presentation
import pandas as pd
import textract
import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings, HuggingFaceInstructEmbeddings
from langchain.vectorstores import FAISS
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from langchain.chat_models import ChatOpenAI
from langchain.llms import HuggingFaceHub
from html_template import css, bot_template, user_template
import json
import openpyxl

def extract_text_from_excel_as_json(file_stream):
    '''Extract text from Excel file stream (.xls or .xlsx) and convert to a nested JSON format wrapped in 'timetable' key'''
    json_output = {}
    
    xls = pd.ExcelFile(file_stream)
    
    # Iterate over each sheet
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name)
        
        # Fill NaN cells that resulted from merged cells
        df.ffill(inplace=True)  # This assumes vertical merging
        
        # Transform the dataframe to the nested structure
        nested_dict = {}
        for _, row in df.iterrows():
            day = row['Unnamed: 0']
            day_schedule = {}
            for time, cell in row.drop('Unnamed: 0').items():
                if pd.notna(cell):
                    details = cell.split(' ')
                    lesson = details[0] + ' ' + details[1] if len(details) > 1 else details[0]
                    venue = details[2] if len(details) > 2 else ''
                    day_schedule[time] = {'lesson': lesson, 'venue': venue}
            nested_dict[day] = day_schedule
        
        json_output[sheet_name] = nested_dict

    # Encapsulating the output within 'timetable' key
    formatted_output = {'timetable': json_output}
        
    return json.dumps(formatted_output, indent=4)

def extract_text_from_excel_as_csv(file_stream):
    '''Extract text from Excel file stream (.xls or .xlsx) and convert to CSV format'''
    csv_texts = []
    
    xls = pd.ExcelFile(file_stream)
    
    # Iterate over each sheet
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name)
        
        # Fill NaN cells that resulted from merged cells
        df.ffill(inplace=True)  # This assumes vertical merging
        
        csv_content = df.to_csv(index=False)  # Convert dataframe to CSV format, without writing row indices
        csv_texts.append(csv_content)
        
    return '\n'.join(csv_texts)  # Join multiple sheets' CSV content with a newline separator

def extract_text_from_excel(file_stream):
    '''Extract text from Excel file stream (.xls or .xlsx)'''
    text = ''
    xls = pd.ExcelFile(file_stream)
    
    # Iterate over each sheet
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name)
        text += sheet_name + '\n'  # Add the sheet name
        text += df.to_string() + '\n\n'  # Convert the sheet data to text
    return text

def extract_text_from_docx(file_stream):
    '''Extract text from DOCX file stream'''
    text = ''
    doc = Document(file_stream)
    for para in doc.paragraphs:
        text += para.text + '\n'
    return text

def extract_text_from_pptx(file_stream):
    '''Extract text from PPTX file stream'''
    text = ''
    prs = Presentation(file_stream)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def extract_text_from_doc(file_stream):
    '''Extract text from DOC file using textract'''
    # Since textract requires file path, we'll save the file temporarily
    temp_path = file_stream.name
    with open(temp_path, "wb") as temp_file:
        temp_file.write(file_stream.getbuffer())
    text = textract.process(temp_path).decode("utf-8")
    return text

def get_pdf_text(file_stream):
    '''Extract text from PDF file stream'''
    text = ''
    pdf_reader = PdfReader(file_stream)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_files(files):
    '''Extract text from files'''
    text = ''
    for file in files:
        try:
            if file.type == 'application/pdf':
                text += get_pdf_text(file)
            elif file.type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                text += extract_text_from_docx(file)
            elif file.type == 'application/msword':  # MIME type for .doc
                text += extract_text_from_doc(file)
            elif file.type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':  # MIME type for .pptx
                text += extract_text_from_pptx(file)
            elif file.type == 'application/vnd.ms-powerpoint':  # MIME type for .ppt
                text += textract.process(file.stream).decode('utf-8')
            elif file.type == 'application/vnd.ms-excel':  # MIME type for .xls
                text += extract_text_from_excel_as_json(file)
            elif file.type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':  # MIME type for .xlsx
                text += extract_text_from_excel_as_json(file)
        except Exception as e:
            print(f"Error processing {file.name}: {e}")
            continue  # If there's an error with one file, move on to the next
    return text

def get_text_chunks(raw_text):
    '''Get text chunks'''
    text_splitter = CharacterTextSplitter(
        separator='\n',
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )

    chunks = text_splitter.split_text(raw_text)
    return chunks

def get_conversation_chain(vector_store):
    '''Get conversation chain'''
    llm = ChatOpenAI()
    # llm = HuggingFaceHub(repo_id='google/flan-t5-xxl', model_kwargs={'temperature':0.5, 'max_length': 512})
    memory = ConversationBufferMemory(
        memory_key='chat_history',
        return_messages=True
    )
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vector_store.as_retriever(),
        memory=memory
    )

    return conversation_chain
    

def get_vector_store(text_chunks):
    '''Get vector store'''
    embeddings = OpenAIEmbeddings()
    # embeddings = HuggingFaceInstructEmbeddings(model_name='hkunlp/instructor-xl')
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    return vector_store

def handle_user_input(user_question):
    '''Handle user input'''
    if st.session_state.conversation is None:
        st.session_state.conversation = get_conversation_chain()
    bot_response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = bot_response['chat_history']

    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace("{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace("{{MSG}}", message.content), unsafe_allow_html=True)
            # st.write(message)


def main():
    load_dotenv()
    st.set_page_config(page_title='Chat with multiple PDFs', page_icon=':shark:', layout='wide')

    st.write(css, unsafe_allow_html=True)

    if 'conversation' not in st.session_state:
        st.session_state.conversation = None

    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = None

    st.title('Chat with multiple PDFs')
    st.header('hat with multiple PDFs')
    user_question = st.text_input("Ask a question about your documents")
    if user_question:
        handle_user_input(user_question)


    with st.sidebar:
        st.subheader('Your documents')
        pdf_docs = st.file_uploader('Upload your documents and click on Process', type=['pdf','docx', 'doc', 'pptx', 'xls', 'xlsx'], accept_multiple_files=True)
        if st.button('Process'):
            with st.spinner('Processing your documents...'):
                # get pdf text
                raw_text = extract_text_from_files(pdf_docs)

                # get the text chunks
                text_chunks = get_text_chunks(raw_text)
                st.write(text_chunks)
                print(text_chunks) 
                # create vector store
                vector_store = get_vector_store(text_chunks)

                # create conversation chain
                st.session_state.conversation =get_conversation_chain(vector_store)
                
    

if __name__ == '__main__':
    main()